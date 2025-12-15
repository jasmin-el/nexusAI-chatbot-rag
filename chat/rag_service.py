"""RAG (Retrieval-Augmented Generation) service for document-based chat.

This module implements the complete RAG pipeline using Azure OpenAI for embeddings
and chat completion, and Pinecone for vector storage and similarity search.

The service handles:
- Document ingestion (PDF and TXT files)
- Text chunking and embedding generation
- Vector storage in Pinecone
- Similarity-based retrieval
- Context-aware response generation
- Conversation title generation
"""

from langchain_openai import AzureChatOpenAI, AzureOpenAIEmbeddings
from langchain_core.messages import HumanMessage, SystemMessage
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader, TextLoader, Docx2txtLoader
from pinecone import Pinecone
import os
from pptx import Presentation
from langchain_core.documents import Document
from dotenv import load_dotenv

load_dotenv()

class RAGService:
    """Retrieval-Augmented Generation service for document-based chat.
    
    This service implements a complete RAG pipeline that combines document retrieval
    with large language model generation. Documents are embedded using Azure OpenAI,
    stored in Pinecone vector database, and retrieved based on semantic similarity
    to user queries.
    
    The service is implemented as a singleton to maintain persistent connections
    to Azure OpenAI and Pinecone, reducing initialization overhead.
    
    Attributes:
        azure_endpoint (str): Azure OpenAI service endpoint URL
        azure_api_key (str): Azure OpenAI API key
        azure_api_version (str): Azure OpenAI API version
        azure_embedding_deployment (str): Name of embedding model deployment
        azure_chat_deployment (str): Name of chat model deployment
        pinecone_api_key (str): Pinecone API key
        index_name (str): Name of Pinecone index
        pinecone_host (str): Pinecone index host URL
        pc (Pinecone): Pinecone client instance
        index (Index): Pinecone index instance
        embeddings (AzureOpenAIEmbeddings): Azure OpenAI embedding model (1024 dims)
        llm (AzureChatOpenAI): Azure OpenAI chat model for response generation
    """
    
    def __init__(self):
        """Initialize RAG service with Azure OpenAI and Pinecone connections.
        
        Loads configuration from environment variables and establishes connections
        to Azure OpenAI (for embeddings and chat) and Pinecone (for vector storage).
        
        Environment Variables Required:
            AZURE_OPENAI_ENDPOINT: Azure OpenAI endpoint URL
            AZURE_OPENAI_API_KEY: Azure OpenAI API key
            AZURE_OPENAI_API_VERSION: API version (default: 2024-02-15-preview)
            AZURE_EMBEDDING_DEPLOYMENT: Embedding model deployment name
            AZURE_CHAT_DEPLOYMENT: Chat model deployment name
            PINECONE_API_KEY: Pinecone API key
            PINECONE_INDEX_NAME: Pinecone index name
            PINECONE_HOST: Pinecone index host URL
        """
        # Azure OpenAI Configuration
        self.azure_endpoint = os.getenv("AZURE_OPENAI_ENDPOINT")
        self.azure_api_key = os.getenv("AZURE_OPENAI_API_KEY")
        self.azure_api_version = os.getenv("AZURE_OPENAI_API_VERSION", "2024-02-15-preview")
        self.azure_embedding_deployment = os.getenv("AZURE_EMBEDDING_DEPLOYMENT")
        self.azure_chat_deployment = os.getenv("AZURE_CHAT_DEPLOYMENT")
        
        # Pinecone Configuration
        self.pinecone_api_key = os.getenv("PINECONE_API_KEY")
        self.index_name = os.getenv("PINECONE_INDEX_NAME")
        self.pinecone_host = os.getenv("PINECONE_HOST")

        # --- Pinecone client (v3 for serverless) ---
        self.pc = Pinecone(api_key=self.pinecone_api_key)
        self.index = self.pc.Index(self.index_name, host=self.pinecone_host)

        # --- Azure OpenAI Embeddings + LLM ---
        self.embeddings = AzureOpenAIEmbeddings(
            azure_endpoint=self.azure_endpoint,
            api_key=self.azure_api_key,
            api_version=self.azure_api_version,
            azure_deployment=self.azure_embedding_deployment,
            dimensions=1024
        )
        
        self.llm = AzureChatOpenAI(
            azure_endpoint=self.azure_endpoint,
            api_key=self.azure_api_key,
            api_version=self.azure_api_version,
            azure_deployment=self.azure_chat_deployment,
            temperature=0.7
        )

    def ingest_file(self, file_path):
        """Process and ingest a document into the RAG knowledge base.
        
        Loads a document, splits it into chunks, generates embeddings, and stores
        them in Pinecone with metadata. Supports PDF and TXT file formats.
        
        The process:
        1. Load document using appropriate loader (PyPDF or TextLoader)
        2. Split into chunks (1000 chars, 200 char overlap)
        3. Generate 1024-dimensional embeddings for each chunk
        4. Create unique IDs using MD5 hash of filename
        5. Store vectors in Pinecone with text and source metadata
        
        Args:
            file_path (str): Absolute path to the document file
            
        Returns:
            str: Success message with number of chunks ingested and filename
            
        Raises:
            ValueError: If file type is not PDF or TXT
            Exception: If document loading, embedding, or Pinecone upsert fails
        """
        if file_path.endswith(".pdf"):
            loader = PyPDFLoader(file_path)
            documents = loader.load()
        elif file_path.endswith(".txt"):
            loader = TextLoader(file_path, encoding="utf-8")
            documents = loader.load()
        elif file_path.endswith(".docx"):
            loader = Docx2txtLoader(file_path)
            documents = loader.load()
        elif file_path.endswith(".pptx"):
            # Custom lightweight loader for PPTX using python-pptx
            prs = Presentation(file_path)
            text_content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
            
            full_text = "\n".join(text_content)
            documents = [Document(page_content=full_text, metadata={"source": file_path})]
        else:
            raise ValueError("Unsupported file type")
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        chunks = splitter.split_documents(documents)

        # Add to Pinecone directly
        import hashlib
        vectors = []
        for i, chunk in enumerate(chunks):
            embedding = self.embeddings.embed_query(chunk.page_content)
            # Create ASCII-safe ID using MD5 hash of filename + index
            file_id = hashlib.md5(os.path.basename(file_path).encode('utf-8')).hexdigest()
            vector_id = f"{file_id}-{i}"
            
            vectors.append({
                "id": vector_id,
                "values": embedding,
                "metadata": {
                    "text": chunk.page_content,
                    "source": os.path.basename(file_path)
                }
            })
        
        self.index.upsert(vectors=vectors)
        return f"Ingested {len(chunks)} chunks from {os.path.basename(file_path)}"

    def generate_response(self, query):
        """Generate AI response using RAG pipeline.
        
        Retrieves relevant document chunks from Pinecone based on semantic similarity
        and uses them as context for Azure OpenAI to generate an informed response.
        
        The process:
        1. Embed the user query using Azure OpenAI
        2. Query Pinecone for top 3 most similar document chunks
        3. Extract text content from retrieved chunks
        4. Build prompt with system message, context, and query
        5. Generate response using Azure OpenAI chat model
        
        Args:
            query (str): User's question or message
            
        Returns:
            str: AI-generated response based on retrieved context
            
        Note:
            Returns "No documents ingested yet." if Pinecone has no vectors
        """
        # Get embedding for query
        query_embedding = self.embeddings.embed_query(query)
        
        # Query Pinecone
        results = self.index.query(
            vector=query_embedding,
            top_k=3,
            include_metadata=True
        )

        if not results['matches']:
            return "No documents ingested yet."

        # Extract context from results
        context = "\n\n".join([match['metadata']['text'] for match in results['matches']])
        
        messages = [
            SystemMessage(content=(
                "Tu es un assistant professionnel et précis. "
                "Réponds de manière claire, structurée et complète, uniquement en texte simple. "
                "Ne mets pas de gras, pas de titres en Markdown, pas d'étoiles, ni de symboles spéciaux. "
                "Chaque fois que tu utilises un tiret pour lister des éléments, commence une nouvelle ligne après le tiret. "
                "Utilise des paragraphes et phrases lisibles pour un style professionnel."
            )),
            HumanMessage(content=f"Contexte:\n{context}\n\nQuestion: {query}")
        ]

        response = self.llm.invoke(messages)
        return response.content

    def generate_title(self, user_message, ai_response):
        """Generate a concise title for a conversation using GPT.
        
        Uses Azure OpenAI to create a short, professional title (5-6 words max)
        based on the first user message and AI response in a conversation.
        
        Args:
            user_message (str): First message from the user
            ai_response (str): AI's response to the first message
            
        Returns:
            str: Generated conversation title (stripped of whitespace)
            
        Note:
            The system prompt instructs the model to generate titles in French,
            without quotes or punctuation.
        """
        messages = [
            SystemMessage(content="Tu es un expert en synthèse. Génère un titre très court (max 5-6 mots) et professionnel pour cette conversation. Ne mets pas de guillemets, pas de point final."),
            HumanMessage(content=f"Message utilisateur: {user_message}\nRéponse AI: {ai_response}\n\nTitre:")
        ]
        response = self.llm.invoke(messages)
        return response.content.strip()

# Singleton instance of RAGService
# Maintains persistent connections to Azure OpenAI and Pinecone
rag_service = RAGService()
